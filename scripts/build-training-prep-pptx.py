#!/usr/bin/env python
"""Build the meeting deck that accompanies 09-teacher-training-prep.md.

The deck is the same findings as the PDF, resequenced for a room: one claim
per slide, the number that supports it large enough to read from the back.

Font note: PowerPoint renders with fonts installed on the machine that opens
the file, so this uses Tahoma — present on every Windows install and correct
for Thai — rather than Sarabun, which would silently fall back elsewhere.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "manual-pdf" / "training-2026-08" / "09-เอกสารเตรียมอบรมครู-สไลด์นำเสนอ.pptx"

FONT = "Tahoma"
NAVY = RGBColor(0x18, 0x3A, 0x5A)
TEAL = RGBColor(0x24, 0x5D, 0x7A)
INK = RGBColor(0x2D, 0x37, 0x48)
MUTED = RGBColor(0x5B, 0x67, 0x70)
GOOD = RGBColor(0x1B, 0x7F, 0x4B)
WARN = RGBColor(0xB4, 0x54, 0x09)
BAD = RGBColor(0xB0, 0x2A, 0x37)
PAPER = RGBColor(0xF6, 0xFA, 0xFC)
LINE = RGBColor(0xD5, 0xE4, 0xEA)


def set_text(frame, runs, *, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
             space_after=6, line=1.35):
    """Fill a text frame from a list of (text, overrides) tuples."""
    frame.word_wrap = True
    # A newline inside a run is not a line break in OOXML — it renders as a
    # stray glyph. Split into real paragraphs instead.
    expanded = []
    for item in runs:
        text, over = item if isinstance(item, tuple) else (item, {})
        lines = text.split("\n")
        for j, part in enumerate(lines):
            o = dict(over)
            if j < len(lines) - 1:
                o["space_after"] = 2  # keep wrapped lines visually together
            expanded.append((part, o))
    for i, item in enumerate(expanded):
        text, over = item
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.alignment = over.get("align", align)
        p.space_after = Pt(over.get("space_after", space_after))
        p.line_spacing = over.get("line", line)
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = FONT
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(slide, x, y, w, h, *, fill=None, line_color=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.06
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def title_slide(prs, eyebrow, title, subtitle, footer):
    s = blank(prs)
    bar = box(s, Cm(0), Cm(0), prs.slide_width, Cm(0.35), fill=NAVY)
    bar.line.fill.background()
    tb = s.shapes.add_textbox(Cm(2.2), Cm(4.4), prs.slide_width - Cm(4.4), Cm(8))
    set_text(tb.text_frame, [
        (eyebrow, {"size": 15, "color": TEAL, "bold": True, "space_after": 10}),
        (title, {"size": 40, "color": NAVY, "bold": True, "line": 1.2, "space_after": 14}),
        (subtitle, {"size": 19, "color": MUTED, "line": 1.4}),
    ])
    fb = s.shapes.add_textbox(Cm(2.2), prs.slide_height - Cm(2.4), prs.slide_width - Cm(4.4), Cm(1.2))
    set_text(fb.text_frame, [(footer, {"size": 13, "color": MUTED})])
    return s


def section_head(prs, slide, title, kicker=None):
    tb = slide.shapes.add_textbox(Cm(1.6), Cm(1.0), prs.slide_width - Cm(3.2), Cm(2.4))
    runs = [(title, {"size": 30, "color": NAVY, "bold": True, "space_after": 4})]
    if kicker:
        runs.append((kicker, {"size": 15, "color": MUTED}))
    set_text(tb.text_frame, runs)


def bullets_slide(prs, title, items, kicker=None):
    s = blank(prs)
    section_head(prs, s, title, kicker)
    tb = s.shapes.add_textbox(Cm(1.9), Cm(4.2), prs.slide_width - Cm(3.8), Cm(11))
    runs = []
    for it in items:
        text, over = it if isinstance(it, tuple) else (it, {})
        over = {"size": 19, "space_after": 13, **over}
        runs.append((text, over))
    set_text(tb.text_frame, runs)
    return s


def stat_slide(prs, title, stats, note=None):
    """Big-number slide. stats is a list of (number, label, color)."""
    s = blank(prs)
    section_head(prs, s, title)
    n = len(stats)
    margin, gap = Cm(1.9), Cm(0.6)
    total_w = prs.slide_width - margin * 2
    w = int((total_w - gap * (n - 1)) / n)
    y, h = Cm(4.6), Cm(6.4)
    for i, (num, label, color) in enumerate(stats):
        x = margin + i * (w + gap)
        box(s, x, y, w, h, fill=PAPER, line_color=LINE)
        tb = s.shapes.add_textbox(x + Cm(0.4), y + Cm(1.0), w - Cm(0.8), h - Cm(1.4))
        set_text(tb.text_frame, [
            (num, {"size": 50, "color": color, "bold": True, "align": PP_ALIGN.CENTER,
                   "space_after": 8, "line": 1.0}),
            (label, {"size": 15, "color": INK, "align": PP_ALIGN.CENTER, "line": 1.3}),
        ])
    if note:
        nb = s.shapes.add_textbox(Cm(1.9), y + h + Cm(0.7), prs.slide_width - Cm(3.8), Cm(2.4))
        set_text(nb.text_frame, [(note, {"size": 15, "color": MUTED, "line": 1.4})])
    return s


def table_slide(prs, title, headers, rows, kicker=None, col_w=None):
    s = blank(prs)
    section_head(prs, s, title, kicker)
    nrows, ncols = len(rows) + 1, len(headers)
    x, y = Cm(1.9), Cm(4.3)
    w = prs.slide_width - Cm(3.8)
    h = Cm(0.95) * nrows
    shape = s.shapes.add_table(nrows, ncols, x, y, w, h)
    tbl = shape.table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = int(w * cw / total)
    for c, head in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.name, p.font.size, p.font.bold = FONT, Pt(15), True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            text, color = val if isinstance(val, tuple) else (val, INK)
            cell = tbl.cell(r, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r % 2 else PAPER
            p = cell.text_frame.paragraphs[0]
            p.font.name, p.font.size = FONT, Pt(14)
            p.font.color.rgb = color
            p.font.bold = color in (GOOD, WARN, BAD)
    return s


def closing_slide(prs, title, questions):
    s = blank(prs)
    section_head(prs, s, title, "ต้องได้ข้อสรุปก่อนประกาศใช้")
    y = Cm(4.4)
    for i, q in enumerate(questions, start=1):
        box(s, Cm(1.9), y, prs.slide_width - Cm(3.8), Cm(2.0), fill=PAPER, line_color=LINE)
        nb = s.shapes.add_textbox(Cm(2.3), y + Cm(0.45), Cm(1.2), Cm(1.2))
        set_text(nb.text_frame, [(str(i), {"size": 24, "color": TEAL, "bold": True})])
        tb = s.shapes.add_textbox(Cm(3.5), y + Cm(0.5), prs.slide_width - Cm(5.8), Cm(1.4))
        set_text(tb.text_frame, [(q, {"size": 17, "color": INK, "line": 1.25})])
        y += Cm(2.25)
    return s


def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Cm(33.87), Cm(19.05)  # 16:9

    title_slide(
        prs,
        "ประชุมเคาะแผนก่อนประกาศใช้ · 27 สิงหาคม 2569",
        "ระบบพร้อมแล้ว\nสิ่งที่ยังไม่พร้อมคือผู้ใช้",
        "ผลทดสอบระบบรถรับส่งนักเรียนจังหวัดลำปาง จากการใช้งานจริงครบทั้ง 6 สิทธิ์",
        "ทุกตัวเลขในเอกสารนี้เก็บจากระบบจริง ไม่ใช่การประมาณการ",
    )

    table_slide(
        prs, "ความพร้อมใน 30 วินาที",
        ["ด้าน", "สถานะ", "ความหมาย"],
        [
            ["ตัวระบบ (ซอฟต์แวร์)", ("พร้อมใช้", GOOD), "ทดสอบครบ 6 สิทธิ์ ไม่พบสิ่งที่ขวางการประกาศใช้"],
            ["การเข้าถึงจากอินเทอร์เน็ต", ("พร้อมใช้", GOOD), "เข้าใช้งานได้จริงทั้งหน้าเว็บและระบบเบื้องหลัง"],
            ["ระบบสิทธิ์", ("พร้อมใช้", GOOD), "กันการเข้าถึงข้ามสิทธิ์ได้จริงที่ฝั่งเซิร์ฟเวอร์"],
            ["ข้อมูลนักเรียน", ("พร้อมใช้", GOOD), "4,462 คน จาก 317 โรงเรียน"],
            ["บัญชีผู้ใช้", ("ต้องเตรียมก่อน", BAD), "643 บัญชียังไม่เคยเข้าระบบ (คนขับ 447 · โรงเรียน 195)"],
            ["ข้อมูลรถ", ("ต้องเตรียมก่อน", BAD), "586 คัน ยังไม่ผ่านการรับรองแม้แต่คันเดียว"],
            ["การเชื่อม LINE ผู้ปกครอง", ("ยังไม่เริ่ม", WARN), "ผูกแล้ว 11 ราย จาก 4,462"],
        ],
        col_w=[26, 20, 54],
    )

    stat_slide(
        prs, "ตัวเลขสามตัวที่อธิบายทุกอย่าง",
        [
            ("643", "บัญชีที่ยังไม่เคยเข้าระบบสำเร็จ\n(คนขับ 447 · โรงเรียน 195)", BAD),
            ("0", "รถที่ผ่านการรับรอง\nจากทั้งหมด 586 คัน", BAD),
            ("11", "ผู้ปกครองที่ผูก LINE แล้ว\nจากทั้งหมด 4,462 ราย", WARN),
        ],
        note="ตัวเลขชุดนี้ไม่ได้แปลว่าระบบเสีย — ทุกอย่างทำงานได้ แต่ยังไม่มีคนเข้ามาใช้ "
             "การเช็กชื่อขึ้น-ลงรถหยุดนิ่งตั้งแต่ 2 กรกฎาคม 2569 ด้วยเหตุผลเดียวกัน",
    )

    bullets_slide(
        prs, "สิ่งที่ทดสอบแล้วว่าใช้ได้จริง",
        [
            ("นำเข้ารายชื่อนักเรียนด้วยไฟล์ CSV", {"bold": True, "color": NAVY, "space_after": 4}),
            "ทดสอบด้วยไฟล์ที่จงใจใส่ข้อผิดพลาด ระบบจับได้ครบ — รหัสซ้ำ ชื่อว่าง เบอร์ผิดรูปแบบ "
            "และให้ตรวจสอบก่อนยืนยันเสมอ ถ้านำเข้าผิดย้อนกลับได้ พร้อมบันทึกเหตุผลลงประวัติ",
            ("ระบบสิทธิ์กันข้อมูลข้ามหน่วยงาน", {"bold": True, "color": NAVY, "space_after": 4}),
            "ทดสอบโดยพยายามเข้าถึงข้อมูลที่ไม่ใช่ของตน ถูกปฏิเสธที่ฝั่งเซิร์ฟเวอร์ทุกครั้ง "
            "ไม่ใช่แค่ซ่อนเมนู — เจ้าหน้าที่ขนส่งไม่เห็นรายชื่อนักเรียน",
            ("การใช้งานบนมือถือของคนขับ", {"bold": True, "color": NAVY, "space_after": 4}),
            "ปุ่มทุกปุ่มไม่ต่ำกว่า 44 พิกเซล กดด้วยนิ้วโป้งได้ · ยังต้องทดสอบเพิ่มเรื่องการอ่านกลางแดด",
        ],
        kicker="ผ่านการกดใช้งานจริง ไม่ใช่การตรวจจากเอกสาร",
    )

    table_slide(
        prs, "คนขับพร้อมแค่ไหน — ไขว้สองตัวเลข",
        ["กลุ่มบัญชีคนขับที่ใช้งานอยู่", "ผูกรถแล้ว", "ยังไม่ผูกรถ", "รวม"],
        [
            ["ยังไม่เคยเข้าระบบ", ("403", GOOD), ("44", WARN), "447"],
            ["เคยเข้าระบบแล้ว", "1", "3", "4"],
            ["รวม", ("404 — ร้อยละ 90", GOOD), ("47", WARN), "451"],
        ],
        kicker="ร้อยละ 90 ผูกรถเรียบร้อยแล้ว — อุปสรรคหลักคือการเข้าระบบครั้งแรก ไม่ใช่การผูกรถ",
        col_w=[40, 20, 20, 20],
    )

    bullets_slide(
        prs, "สามอย่างที่ต้องเตรียมก่อนวันอบรม",
        [
            ("1 · ผูกรถให้คนขับ 44 รายที่เหลือ", {"bold": True, "color": WARN, "space_after": 4}),
            "ข่าวดี: คนขับ 403 จาก 447 รายที่ยังไม่เคยเข้าระบบ ผูกรถเรียบร้อยแล้ว พร้อมใช้ทันทีที่เข้าได้ "
            "เหลือ 44 รายที่ต้องผูกก่อน มิฉะนั้นเข้าไปแล้วใช้งานอะไรไม่ได้เลย",
            ("2 · วางแผนพาผู้ใช้เข้าระบบครั้งแรก", {"bold": True, "color": BAD, "space_after": 4}),
            "643 บัญชีได้รับรหัสไปแล้วแต่ยังไม่เคยเข้า การอบรมต้องเผื่อเวลาสำหรับ "
            "ขั้นตอนเปลี่ยนรหัสผ่านครั้งแรกของทุกคน",
            ("3 · ตัดสินใจเรื่องกระบวนการรับรองรถ", {"bold": True, "color": BAD, "space_after": 4}),
            "586 คันยังไม่ผ่านการรับรอง ข้อมูลทะเบียน พ.ร.บ. ภาษี ยังไม่มีในระบบ "
            "ต้องกำหนดว่าใครกรอก และเสร็จภายในกี่เดือน",
        ],
        kicker="ถ้าข้ามหัวข้อนี้ การอบรมจะติดขัดตั้งแต่ต้นชั่วโมง",
    )

    table_slide(
        prs, "ลำดับการอบรมที่แนะนำ",
        ["ช่วง", "เนื้อหา", "หมายเหตุ"],
        [
            ["1", "เข้าสู่ระบบ + เปลี่ยนรหัสครั้งแรก", ("เผื่อเวลาให้มาก คนส่วนใหญ่ยังไม่เคยเข้า", WARN)],
            ["2", "ดูแดชบอร์ดและรายชื่อนักเรียนของโรงเรียนตนเอง", "สร้างความคุ้นเคยก่อนลงมือแก้ข้อมูล"],
            ["3", "นำเข้ารายชื่อนักเรียนด้วยไฟล์", ("ใช้ไฟล์ทดสอบ อย่าใช้ไฟล์จริง", BAD)],
            ["4", "ตรวจผลก่อนยืนยัน และการย้อนกลับ", "จุดที่ครูกังวลมากที่สุด ควรให้ลองเอง"],
            ["5", "จัดการรถและจุดรับส่ง", ""],
            ["6", "ถาม-ตอบ และแจกช่องทางขอความช่วยเหลือ", ""],
        ],
        col_w=[8, 48, 44],
    )

    bullets_slide(
        prs, "ข้อจำกัดที่ควรบอกตั้งแต่ต้น",
        [
            "การเชื่อม LINE ผู้ปกครองยังไม่แพร่หลาย — ผูกแล้ว 11 รายจาก 4,462 "
            "ยังไม่ควรประกาศกับผู้ปกครองว่าใช้ได้แล้ว",
            "การเช็กชื่อขึ้น-ลงรถหยุดนิ่งตั้งแต่ 2 กรกฎาคม 2569 — ไม่ใช่ระบบเสีย แต่ยังไม่มีคนขับใช้งาน",
            "จุดรับส่งมีเพียง 27 จุด จากรถ 586 คัน — ข้อมูลเส้นทางยังไม่ครบ",
            "หน้าผู้ปกครองทดสอบนอกแอป LINE ไม่ได้ — เป็นข้อจำกัดโดยการออกแบบเพื่อความปลอดภัย",
            "การอ่านหน้าจอกลางแดดและการกดขณะรถสั่นยังไม่ได้ทดสอบ — ต้องลองบนมือถือจริง",
        ],
        kicker="การบอกข้อจำกัดตั้งแต่ต้น ดีกว่าให้ครูไปเจอเอง",
    )

    closing_slide(
        prs, "สี่คำถามที่ต้องเคาะวันนี้",
        [
            "เป้าหมายจำนวนบัญชีที่ต้องเข้าระบบสำเร็จ ภายในกี่วันหลังประกาศใช้",
            "ใครรับผิดชอบผูกรถให้คนขับ 44 รายที่เหลือ และเสร็จเมื่อไร",
            "กระบวนการตรวจรับรองรถจะเริ่มเมื่อไร ตั้งเป้ากี่คันต่อเดือน",
            "จะประกาศให้ผู้ปกครองใช้ LINE เมื่อผูกบัญชีได้กี่เปอร์เซ็นต์",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    print(build())
